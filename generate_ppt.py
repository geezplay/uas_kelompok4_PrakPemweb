import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    DARK_BG = RGBColor(10, 10, 10)       # Matte Black
    DARK_CARD = RGBColor(24, 24, 24)     # Dark Gray Card
    LIGHT_BG = RGBColor(250, 250, 250)   # Light Gray
    LIGHT_CARD = RGBColor(255, 255, 255) # Pure White
    RED_ACCENT = RGBColor(220, 38, 38)   # Racing Red
    GREEN_ACCENT = RGBColor(16, 185, 129) # Success Emerald
    TEXT_MUTED = RGBColor(156, 163, 175) # Gray-400
    TEXT_LIGHT = RGBColor(243, 244, 246) # Gray-100
    TEXT_DARK = RGBColor(17, 24, 39)     # Gray-900
    
    # Helper to add background color
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to add text boxes
    def add_text_box(slide, left, top, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT, font_name="Arial"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        return tf

    # Helper to add shape card
    def add_card(slide, left, top, width, height, bg_color, border_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    blank_slide_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: TITLE SLIDE (Dark Mode)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1, DARK_BG)
    
    # Accent Line (Top red bar)
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RED_ACCENT
    accent_bar.line.fill.background()

    # Large BG graphic card
    add_card(slide1, Inches(1), Inches(1.5), Inches(11.333), Inches(4.5), DARK_CARD)

    # Category Tag
    add_text_box(slide1, Inches(1.5), Inches(2.0), Inches(8), Inches(0.4), "🏎️ TUGAS AKHIR • PEMROGRAMAN WEB", 14, RED_ACCENT, bold=True)
    
    # Title
    add_text_box(slide1, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2), "RACE SETUP GARAGE", 54, TEXT_LIGHT, bold=True, font_name="Arial Black")
    
    # Subtitle
    add_text_box(slide1, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8), "Sistem Manajemen Setup Motor Balap Berbasis Web CRUD\nImplementasi React & Tailwind CSS Tanpa Database", 18, TEXT_MUTED)

    # Presenter Details Card
    presenter_y = Inches(4.8)
    add_text_box(slide1, Inches(1.5), presenter_y, Inches(10), Inches(0.4), "Disusun Untuk Memenuhi Nilai Ujian Akhir Semester (UAS)", 13, TEXT_MUTED, bold=True)
    add_text_box(slide1, Inches(1.5), presenter_y + Inches(0.3), Inches(10), Inches(0.4), "Kelompok Pemrograman Web • Deployment: Vercel", 12, TEXT_MUTED)


    # ==========================================
    # SLIDE 2: LATAR BELAKANG (Light Mode)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2, LIGHT_BG)

    # Slide Title
    add_text_box(slide2, Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.5), "LATAR BELAKANG", 14, RED_ACCENT, bold=True)
    add_text_box(slide2, Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8), "Mengapa Race Setup Garage Diperlukan?", 32, TEXT_DARK, bold=True)

    # Left Card (The Problem)
    add_card(slide2, Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.3), LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide2, Inches(1.4), Inches(2.5), Inches(4.5), Inches(0.4), "⚠️ MASALAH UTAMA", 18, RED_ACCENT, bold=True)
    problem_text = (
        "• Catatan Manual: Mekanik balap sering kali mencatat setup motor (Main Jet, Pilot Jet, Gir) secara manual di kertas.\n\n"
        "• Risiko Tinggi: Kertas catatan mudah robek, basah karena oli/air, atau hilang di area paddock yang sibuk.\n\n"
        "• Sulit Dianalisis: Data historis setelan motor sulit dicari dengan cepat ketika kondisi cuaca sirkuit berubah tiba-tiba."
    )
    add_text_box(slide2, Inches(1.4), Inches(3.1), Inches(4.5), Inches(3.0), problem_text, 14, TEXT_DARK)

    # Right Card (The Solution)
    add_card(slide2, Inches(7.0), Inches(2.2), Inches(5.3), Inches(4.3), LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide2, Inches(7.4), Inches(2.5), Inches(4.5), Inches(0.4), "🏁 SOLUSI DIGITAL", 18, GREEN_ACCENT, bold=True)
    solution_text = (
        "• Race Setup Garage: Aplikasi pencatatan digital berbasis web khusus untuk setelan motor balap.\n\n"
        "• Tanpa Database Eksternal: Memanfaatkan LocalStorage client-side yang andal, memastikan data tetap tersimpan secara instan tanpa biaya server.\n\n"
        "• Spec Sheet Siap Cetak: Ekspor PDF & Cetak lembar resmi mekanik di sirkuit secara teratur."
    )
    add_text_box(slide2, Inches(7.4), Inches(3.1), Inches(4.5), Inches(3.0), solution_text, 14, TEXT_DARK)


    # ==========================================
    # SLIDE 3: TUJUAN PEMBUATAN APLIKASI (Light Mode)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3, LIGHT_BG)

    add_text_box(slide3, Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.5), "TUJUAN APLIKASI", 14, RED_ACCENT, bold=True)
    add_text_box(slide3, Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8), "Fokus & Manfaat Utama Sistem", 32, TEXT_DARK, bold=True)

    # 3 Column Cards
    card_width = Inches(3.5)
    card_height = Inches(4.3)
    gap = Inches(0.4)
    start_x = Inches(1.0)
    y_pos = Inches(2.2)

    # Card 1: Efisiensi
    add_card(slide3, start_x, y_pos, card_width, card_height, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide3, start_x + Inches(0.3), y_pos + Inches(0.4), card_width - Inches(0.6), Inches(0.4), "01. Efisiensi Kerja", 18, RED_ACCENT, bold=True)
    text_1 = "Mempercepat waktu mekanik dalam mencari setelan karburator dan gir rasio terbaik yang pernah dicoba pada sirkuit tertentu sebelumnya."
    add_text_box(slide3, start_x + Inches(0.3), y_pos + Inches(1.0), card_width - Inches(0.6), Inches(2.5), text_1, 14, TEXT_DARK)

    # Card 2: Portabilitas
    x2 = start_x + card_width + gap
    add_card(slide3, x2, y_pos, card_width, card_height, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide3, x2 + Inches(0.3), y_pos + Inches(0.4), card_width - Inches(0.6), Inches(0.4), "02. Portabilitas Tinggi", 18, RED_ACCENT, bold=True)
    text_2 = "Aplikasi dapat diakses dari smartphone, tablet, atau laptop mekanik langsung di paddock tanpa ketergantungan koneksi database yang lambat."
    add_text_box(slide3, x2 + Inches(0.3), y_pos + Inches(1.0), card_width - Inches(0.6), Inches(2.5), text_2, 14, TEXT_DARK)

    # Card 3: Standardisasi
    x3 = x2 + card_width + gap
    add_card(slide3, x3, y_pos, card_width, card_height, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide3, x3 + Inches(0.3), y_pos + Inches(0.4), card_width - Inches(0.6), Inches(0.4), "03. Standardisasi Data", 18, RED_ACCENT, bold=True)
    text_3 = "Menyediakan form standar yang lengkap (Jetting, Needle Clip, Ignition, Gear, Temperatur, Cuaca) untuk meminimalkan salah komunikasi antar mekanik."
    add_text_box(slide3, x3 + Inches(0.3), y_pos + Inches(1.0), card_width - Inches(0.6), Inches(2.5), text_3, 14, TEXT_DARK)


    # ==========================================
    # SLIDE 4: ALUR KERJA WEBSITE (Light Mode) - NEW DETAIL SLIDE
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4, LIGHT_BG)

    add_text_box(slide4, Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.5), "ALUR KERJA PENGGUNA", 14, RED_ACCENT, bold=True)
    add_text_box(slide4, Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8), "Bagaimana Mekanik Menggunakan Aplikasi?", 32, TEXT_DARK, bold=True)

    # 4 horizontal steps mapping the workflow
    step_w = Inches(2.6)
    step_h = Inches(4.3)
    step_gap = Inches(0.3)
    step_start_x = Inches(1.0)
    step_y = Inches(2.2)

    # Step 1: Monitor
    add_card(slide4, step_start_x, step_y, step_w, step_h, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide4, step_start_x + Inches(0.2), step_y + Inches(0.3), step_w - Inches(0.4), Inches(0.4), "Langkah 01", 12, RED_ACCENT, bold=True)
    add_text_box(slide4, step_start_x + Inches(0.2), step_y + Inches(0.7), step_w - Inches(0.4), Inches(0.5), "Monitoring & Analisis", 15, TEXT_DARK, bold=True)
    text_s1 = "Mekanik membuka aplikasi dan memantau rangkuman statistik di Dashboard (total setup, motor terpopuler, rerata suhu sirkuit saat ini)."
    add_text_box(slide4, step_start_x + Inches(0.2), step_y + Inches(1.4), step_w - Inches(0.4), Inches(2.5), text_s1, 12, TEXT_DARK)

    # Step 2: Search & Filter
    x_s2 = step_start_x + step_w + step_gap
    add_card(slide4, x_s2, step_y, step_w, step_h, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide4, x_s2 + Inches(0.2), step_y + Inches(0.3), step_w - Inches(0.4), Inches(0.4), "Langkah 02", 12, RED_ACCENT, bold=True)
    add_text_box(slide4, x_s2 + Inches(0.2), step_y + Inches(0.7), step_w - Inches(0.4), Inches(0.5), "Pencarian & Filter", 15, TEXT_DARK, bold=True)
    text_s2 = "Mencari kata kunci atau menggunakan filter sirkuit & cuaca untuk menemukan setelan karburator/gir yang cocok dengan kondisi lintasan saat itu."
    add_text_box(slide4, x_s2 + Inches(0.2), step_y + Inches(1.4), step_w - Inches(0.4), Inches(2.5), text_s2, 12, TEXT_DARK)

    # Step 3: CRUD
    x_s3 = x_s2 + step_w + step_gap
    add_card(slide4, x_s3, step_y, step_w, step_h, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide4, x_s3 + Inches(0.2), step_y + Inches(0.3), step_w - Inches(0.4), Inches(0.4), "Langkah 03", 12, RED_ACCENT, bold=True)
    add_text_box(slide4, x_s3 + Inches(0.2), step_y + Inches(0.7), step_w - Inches(0.4), Inches(0.5), "Manajemen Data (CRUD)", 15, TEXT_DARK, bold=True)
    text_s3 = "Mekanik menambahkan setup baru hasil penyesuaian, mengubah spesifikasi yang kurang optimal, atau menghapus setup yang salah via modal form."
    add_text_box(slide4, x_s3 + Inches(0.2), step_y + Inches(1.4), step_w - Inches(0.4), Inches(2.5), text_s3, 12, TEXT_DARK)

    # Step 4: Print/Export
    x_s4 = x_s3 + step_w + step_gap
    add_card(slide4, x_s4, step_y, step_w, step_h, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide4, x_s4 + Inches(0.2), step_y + Inches(0.3), step_w - Inches(0.4), Inches(0.4), "Langkah 04", 12, RED_ACCENT, bold=True)
    add_text_box(slide4, x_s4 + Inches(0.2), step_y + Inches(0.7), step_w - Inches(0.4), Inches(0.5), "Output & Distribusi", 15, TEXT_DARK, bold=True)
    text_s4 = "Mengekspor data ke file CSV (Excel) untuk cadangan, atau mencetak Spec Sheet resmi ke PDF untuk ditempel di dashboard motor balap."
    add_text_box(slide4, x_s4 + Inches(0.2), step_y + Inches(1.4), step_w - Inches(0.4), Inches(2.5), text_s4, 12, TEXT_DARK)


    # ==========================================
    # SLIDE 5: ALIRAN DATA / SISTEM (Dark Mode)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5, DARK_BG)

    add_text_box(slide5, Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.5), "ALUR KERJA SISTEM", 14, RED_ACCENT, bold=True)
    add_text_box(slide5, Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8), "Aliran Sinkronisasi Data Client-Side", 32, TEXT_LIGHT, bold=True)

    # Left Side: Explanation Box
    add_card(slide5, Inches(1.0), Inches(2.2), Inches(4.5), Inches(4.3), DARK_CARD)
    add_text_box(slide5, Inches(1.3), Inches(2.5), Inches(3.9), Inches(0.4), "Alur Sinkronisasi State", 18, RED_ACCENT, bold=True)
    architecture_text = (
        "• Single-Page React App: State diatur terpusat di root component (page.js) dan dialirkan sebagai props ke sub-komponen.\n\n"
        "• LocalStorage Hydration: Saat aplikasi dimuat, data di-load dari browser client. Saat data diubah (CRUD), local storage langsung sinkron secara otomatis.\n\n"
        "• SSR-Friendly: Logika inisialisasi diletakkan di useEffect untuk mencegah error window is not defined."
    )
    add_text_box(slide5, Inches(1.3), Inches(3.0), Inches(3.9), Inches(3.2), architecture_text, 13, TEXT_LIGHT)

    # Right Side: Visual Flow Chart (using text cards to simulate blocks)
    flow_x = Inches(6.0)
    flow_y = Inches(2.2)
    flow_w = Inches(6.333)
    
    # Step 1
    add_card(slide5, flow_x, flow_y, flow_w, Inches(1.0), DARK_CARD, border_color=RED_ACCENT)
    add_text_box(slide5, flow_x + Inches(0.3), flow_y + Inches(0.15), flow_w - Inches(0.6), Inches(0.7), "User Interface (React Components)\nRender data dan menangani event aksi dari Mekanik (Tambah, Edit, Hapus).", 12, TEXT_LIGHT, bold=True)
    
    # Arrow 1
    add_text_box(slide5, flow_x, flow_y + Inches(1.05), flow_w, Inches(0.3), "▼", 14, RED_ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # Step 2
    flow_y2 = flow_y + Inches(1.4)
    add_card(slide5, flow_x, flow_y2, flow_w, Inches(1.0), DARK_CARD, border_color=RED_ACCENT)
    add_text_box(slide5, flow_x + Inches(0.3), flow_y2 + Inches(0.15), flow_w - Inches(0.6), Inches(0.7), "React State Management (page.js)\nMelakukan pemrosesan CRUD, pencarian, filter sirkuit/cuaca, dan pengurutan.", 12, TEXT_LIGHT, bold=True)

    # Arrow 2
    add_text_box(slide5, flow_x, flow_y2 + Inches(1.05), flow_w, Inches(0.3), "▼", 14, RED_ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # Step 3
    flow_y3 = flow_y2 + Inches(1.4)
    add_card(slide5, flow_x, flow_y3, flow_w, Inches(1.0), DARK_CARD, border_color=RED_ACCENT)
    add_text_box(slide5, flow_x + Inches(0.3), flow_y3 + Inches(0.15), flow_w - Inches(0.6), Inches(0.7), "Client-Side Storage (LocalStorage Stream)\nMenyimpan data balap secara persisten dalam format JSON stream di browser lokal.", 12, TEXT_LIGHT, bold=True)


    # ==========================================
    # SLIDE 6: SPESIFIKASI TEKNOLOGI (Light Mode)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6, LIGHT_BG)

    add_text_box(slide6, Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.5), "TEKNOLOGI & TOOLS", 14, RED_ACCENT, bold=True)
    add_text_box(slide6, Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8), "Tech Stack yang Digunakan", 32, TEXT_DARK, bold=True)

    # 4 Grid Tech Cards
    grid_w = Inches(5.3)
    grid_h = Inches(2.0)
    
    # Row 1 Left: Next.js
    add_card(slide6, Inches(1.0), Inches(2.2), grid_w, grid_h, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide6, Inches(1.3), Inches(2.4), grid_w - Inches(0.6), Inches(0.3), "⚛️ React & Next.js 15 (App Router)", 16, TEXT_DARK, bold=True)
    add_text_box(slide6, Inches(1.3), Inches(2.8), grid_w - Inches(0.6), Inches(1.2), "Menggunakan Next.js 15 sebagai engine dasar. Struktur kode modular berbasis App Router (page.js, layout.js) mempermudah manajemen dependensi.", 12, TEXT_MUTED)

    # Row 1 Right: Tailwind CSS
    add_card(slide6, Inches(7.0), Inches(2.2), grid_w, grid_h, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide6, Inches(7.3), Inches(2.4), grid_w - Inches(0.6), Inches(0.3), "🎨 Tailwind CSS v3", 16, TEXT_DARK, bold=True)
    add_text_box(slide6, Inches(7.3), Inches(2.8), grid_w - Inches(0.6), Inches(1.2), "Menyediakan utility-first styling yang responsif. Tampilan mendukung tema transisi Light/Dark Mode yang modern serta animasi micro-interactions.", 12, TEXT_MUTED)

    # Row 2 Left: LocalStorage
    add_card(slide6, Inches(1.0), Inches(4.5), grid_w, grid_h, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide6, Inches(1.3), Inches(4.7), grid_w - Inches(0.6), Inches(0.3), "💾 LocalStorage (JSON Stream)", 16, TEXT_DARK, bold=True)
    add_text_box(slide6, Inches(1.3), Inches(5.1), grid_w - Inches(0.6), Inches(1.2), "Penyimpanan client-side bebas database backend. Kecepatan baca/tulis data instan, bebas kuota, dan data tetap persisten meskipun halaman di-refresh.", 12, TEXT_MUTED)

    # Row 2 Right: Vercel & Git
    add_card(slide6, Inches(7.0), Inches(4.5), grid_w, grid_h, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide6, Inches(7.3), Inches(4.7), grid_w - Inches(0.6), Inches(0.3), "🚀 Deployment & Version Control", 16, TEXT_DARK, bold=True)
    add_text_box(slide6, Inches(7.3), Inches(5.1), grid_w - Inches(0.6), Inches(1.2), "Kode disimpan di GitHub dan di-deploy ke Vercel untuk pembaruan cepat (CI/CD). Memenuhi standar tugas akhir yang dapat diakses publik.", 12, TEXT_MUTED)


    # ==========================================
    # SLIDE 7: FITUR APLIKASI (Light Mode)
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7, LIGHT_BG)

    add_text_box(slide7, Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.5), "FITUR APLIKASI", 14, RED_ACCENT, bold=True)
    add_text_box(slide7, Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8), "Kapabilitas Sistem Manajemen Balap", 32, TEXT_DARK, bold=True)

    # 4 grid features
    feat_w = Inches(5.3)
    feat_h = Inches(2.0)

    # Feat 1: CRUD
    add_card(slide7, Inches(1.0), Inches(2.2), feat_w, feat_h, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide7, Inches(1.3), Inches(2.4), feat_w - Inches(0.6), Inches(0.3), "⚡ Operasi CRUD Penuh", 16, RED_ACCENT, bold=True)
    add_text_box(slide7, Inches(1.3), Inches(2.8), feat_w - Inches(0.6), Inches(1.2), "Menambah setup motor balap baru, melihat lembar spesifikasi terperinci, mengedit detail (Jetting, Ratio), dan menghapus setup dengan konfirmasi modal.", 12, TEXT_DARK)

    # Feat 2: Dashboard Stats
    add_card(slide7, Inches(7.0), Inches(2.2), feat_w, feat_h, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide7, Inches(7.3), Inches(2.4), feat_w - Inches(0.6), Inches(0.3), "📊 Dashboard & Real-Time Stats", 16, RED_ACCENT, bold=True)
    add_text_box(slide7, Inches(7.3), Inches(2.8), feat_w - Inches(0.6), Inches(1.2), "Menghitung total data setup, motor terpopuler, sirkuit teraktif, dan rata-rata suhu lintasan secara dinamis langsung dari dataset aktif.", 12, TEXT_DARK)

    # Feat 3: Filter & Sort
    add_card(slide7, Inches(1.0), Inches(4.5), feat_w, feat_h, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide7, Inches(1.3), Inches(4.7), feat_w - Inches(0.6), Inches(0.3), "🔍 Pencarian & Filter Canggih", 16, RED_ACCENT, bold=True)
    add_text_box(slide7, Inches(1.3), Inches(5.1), feat_w - Inches(0.6), Inches(1.2), "Pencarian teks instan (motor, mekanik, catatan), filter drop-down berdasarkan sirkuit dan cuaca, serta penyortiran (terbaru, suhu, abjad).", 12, TEXT_DARK)

    # Feat 4: Export & Print
    add_card(slide7, Inches(7.0), Inches(4.5), feat_w, feat_h, LIGHT_CARD, border_color=RGBColor(229, 231, 235))
    add_text_box(slide7, Inches(7.3), Inches(4.7), feat_w - Inches(0.6), Inches(0.3), "🖨️ Ekspor Excel (CSV) & Cetak Spec Sheet (PDF)", 16, RED_ACCENT, bold=True)
    add_text_box(slide7, Inches(7.3), Inches(5.1), feat_w - Inches(0.6), Inches(1.2), "Mengunduh semua data setup ke berkas CSV dengan sekali klik. Cetak lembar spesifikasi resmi mekanik (Spec Sheet) yang rapi langsung ke berkas PDF.", 12, TEXT_DARK)


    # ==========================================
    # SLIDE 8: PENUTUP (Dark Mode)
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8, DARK_BG)

    # Accent Line (Top red bar)
    accent_bar8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    accent_bar8.fill.solid()
    accent_bar8.fill.fore_color.rgb = RED_ACCENT
    accent_bar8.line.fill.background()

    # Large Card
    add_card(slide8, Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5), DARK_CARD)

    # Muted Tag
    add_text_box(slide8, Inches(1.5), Inches(2.2), Inches(10), Inches(0.4), "RACE SETUP GARAGE", 14, RED_ACCENT, bold=True)
    
    # Big Thank You
    add_text_box(slide8, Inches(1.5), Inches(2.7), Inches(10), Inches(1.2), "SEKIAN & TERIMA KASIH", 44, TEXT_LIGHT, bold=True, font_name="Arial Black")
    
    # Subinfo
    add_text_box(slide8, Inches(1.5), Inches(4.0), Inches(10), Inches(0.8), "Mari lanjut ke sesi Demonstrasi Aplikasi (CRUD, Ekspor CSV, Cetak PDF)\ndan Tinjauan Repository GitHub.", 16, TEXT_MUTED)

    prs.save("Race_Setup_Garage_Presentasi_V2.pptx")
    print("Presentation generated successfully with User & System Workflow details!")

if __name__ == "__main__":
    create_presentation()
